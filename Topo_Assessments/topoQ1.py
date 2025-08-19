import pandas as pd
import numpy as np
import json
import math
from IPython.display import display
import pdfplumber
from pptx import Presentation

'''
    The code comprises a series of class instances that are nested into one another: 

    industry class (Superclass)
     - Industry Quarter performance class
     - company class 
        - clients class
        - employee class
        - Annual Performance class
           - Quarter Performance class

    The Importer class imports the 4 datasets and uses the data to construct other class instances which are nested. The
    nested structure is then converted into a unified dictionary structure and exported as a json file in local directory.

    Each class except Importer class has a get_dict() function that reformats the instance into a subset of the 
    unified dictionary, before being merged with the dictionary of it's parent class.
    '''
#===========================================================================================================================
# Class List 

    
class Industry:

    def __init__(self, name, industry_Performance_List, company_List):
        self.name                       = name
        self.industry_Performance_List  = industry_Performance_List # 2D array of: [Year] BY [Quarter, Revenue, Memberships sold, AVG Duration(Minutes)]
        self.company_List               = company_List              # 1D array of Company objects

    def get_dict(self):
        json_Company_List = []
        json_Industry_performance_list = []

        for company in self.company_List:
            json_Company_List.append(company.get_dict())

        for quarter in self.industry_Performance_List:

            json_quarter = {"year": quarter[0],
                            "quarter": quarter[1],
                            "revenue": quarter[2],
                            "memberships_sold": quarter[3],
                            "avg_duration_mins": quarter[4]}
            
            json_Industry_performance_list.append(json_quarter)

        dict_obj = {
            "name": self.name,
            "industry_Performance_List": json_Industry_performance_list,
            "company_List": json_Company_List
        }
        return dict_obj

class Company:

    def __init__(self, id, name,  employee_List, quarter_performance_list, industry = "Nan", company_Revenue = "NaN",  location = "NaN", client_List = "NaN"):
        self.id                         = id
        self.name                       = name
        self.industry                   = industry
        self.company_Revenue            = company_Revenue
        self.location                   = location
        self.employee_List              = employee_List             # 1D array of Employee objects
        self.quarter_performance_list    = quarter_performance_list # 1D array of QuarterPerformance objects
        self.client_List                = []                        # 1D array of client objects
        self.annual_performance_list    = []                        # 1D array of AnnualPerformance objects

    def client_setter(self, client_List):
        self.client_List = client_List

    def annual_performance_add(self, val):
        self.annual_performance_list.append(val)
    
    def get_dict(self):
        json_employee_list           = []
        json_quarter_performance_list= []
        json_client_list             = []
        json_annual_performance_list = []

        for employee in self.employee_List:
            json_employee_list.append(employee.get_dict())

        for quarter in self.annual_performance_list:
            json_quarter_performance_list.append(quarter.get_dict())
        
        for client in self.client_List:
            json_client_list.append(client.get_dict())

        for annual in self.annual_performance_list:
            json_annual_performance_list.append(annual.get_dict())

        dict_obj = {
            "id": self.id,
            "name": self.name,
            "industry": self.industry,
            "company_Revenue": self.company_Revenue,
            "location": self.location,
            "employee_List": json_employee_list,
            "quarter_performance_list": json_quarter_performance_list,
            "client_List":  json_client_list,
            "annual_performance_list": json_annual_performance_list
        }
        return dict_obj

class Client:

    def __init__(self, date, membership_id, membership_type, activity, revenue, duration_min, location ):
        self.date                       = date
        self.membership_id              = membership_id
        self.membership_type            = membership_type
        self.activity                   = activity
        self.revenue                    = revenue
        self.duration_min               = duration_min
        self.location                   = location

    def get_dict(self):
        dict_obj = {
            "date": self.date,
            "membership_id": self.membership_id,
            "activity": self.activity,
            "revenue": self.revenue,
            "duration_min": self.duration_min,
            "location": self.location
        }
        return dict_obj

class Employee:

    def __init__(self, id, name, role = "NaN", cashmoney = "NaN", hired_Date = "NaN"):
        self.id                         = id
        self.name                       = name
        self.role                       = role
        self.cashmoney                  = cashmoney
        self.hired_Date                 = hired_Date

    def get_dict(self):
        dict_obj = {
            "id": self.id,
            "name": self.name,
            "role": self.role,
            "cashmoney": self.cashmoney,
            "hired_Date": self.hired_Date
        }
        return dict_obj
    
class QuarterPerformance: 

    def __init__(self, year = "Nan", quarter = "NaN", memberships_sold = "NaN", avg_duration_min = "NaN", revenue = "NaN", profit_margin = "NaN"):
        self.year                       = year
        self.quarter                    = quarter           # This should be a string for compatibility with different documents
        self.revenue                    = revenue
        self.memberships_sold           = memberships_sold
        self.avg_duration_min           = avg_duration_min
        self.profit_margin              = profit_margin

    def get_dict(self):
        dict_obj = {
            "year": self.year,
            "quarter": self.quarter,
            "memberships_sold": self.memberships_sold,
            "avg_duration_min": self.avg_duration_min,
            "revenue": self.revenue,
            "profit_margin": self.profit_margin
        }
        return dict_obj

class AnnualPerformance:

    def __init__(self, year, total_revenue, total_membership_sold, top_location, quarters, revenue_distribution):
        self.year                       = year
        self.total_revenue              = total_revenue
        self.total_membership_sold      = total_membership_sold
        self.top_location               = top_location
        self.quarters                   = quarters              # 1D array of QuarterPerformance objects
        self.revenue_distribution       = revenue_distribution  # 1D arrary of [gym, pool, tennis_Court, personal_Training]

    def get_dict(self):
        json_revenue_distribution = {"gym": self.revenue_distribution[0],
                                     "pool": self.revenue_distribution[1],
                                     "tennis_Court": self.revenue_distribution[2],
                                     "personal_Training": self.revenue_distribution[3]}
        
        json_quarter_list = []

        for quarter in self.quarters:
            json_quarter_list.append(quarter.get_dict())

        dict_obj = {
            "year": self.year,
            "total_revenue": self.total_revenue,
            "total_membership_sold": self.total_membership_sold,
            "top_location": self.top_location,
            "quarters": json_quarter_list,
            "revenue_distribution": json_revenue_distribution
        }
        return dict_obj

class Importer:

    def __init__(self, file_JSON, file_CSV, file_PDF, file_PPTX):
        self.file_json = file_JSON
        self.all_employee               = [] # temp 1D array of employee objects 
        self.all_quarter_performance    = [] # temp 1D array of performance objects
        self.all_companies              = [] # temp 1D array of company objects
        self.all_industry               = [] # 1D array of all industry

        self.import_json(file_JSON)
        self.import_csv(file_CSV)
        self.import_pdf(file_PDF)
        self.import_pptx(file_PPTX)


    def import_json(self,  file):
        # Step 1: Import file as a dataframe
        df_json = pd.read_json(file)
        df_normalized = pd.json_normalize(
            df_json['companies']
        )

        # Step 2: Break down dataframe into smaller datasets, then construct objects to store said data
        # Step 2.1: Break down employee data by company
        for index, company in df_normalized.iterrows():
            employee_List = []
            for employees in company.employees:
                employee = []
                for attribute in employees.values():
                    employee.append(attribute)
                try:
                    employee_temp = Employee(employee[0],
                                    employee[1],
                                    employee[2],
                                    employee[3],
                                    employee[4])
                except:

                    employee_temp = Employee(employee[0],
                                    employee[1],
                                    employee[2],
                                    employee[3])
                employee_List.append(employee_temp)
            if len(employee_List) > 0:
                self.all_employee.append(employee_List)
                
        # Step 2.2: Break down performance data by company
        # Had to reimport the json.csv again as a dict because the original dataframe was hard to manipulate for this stage
        with open('dataset1.json', 'r') as file:
            dict_json = json.load(file)
        for company in range(len(dict_json['companies'])):
            quarter_performance_tmp = dict_json['companies'][company]["performance"]    # output is a 2 layer nested dictionary
            quarter_performance_list = []
            for quarter in quarter_performance_tmp:
                

                quarter_performance_tmp_obj = QuarterPerformance( "NaN",
                                                   quarter,
                                                   "NaN",
                                                   "NaN",
                                                    quarter_performance_tmp[quarter]['revenue'],
                                                    quarter_performance_tmp[quarter]['profit_margin'])
                quarter_performance_list.append(quarter_performance_tmp_obj)
            self.all_quarter_performance.append(quarter_performance_list)

        # Step 2.3: Break down company info data for each company
        selected_range = df_normalized.loc[:, 'id':'location'] 
        for index, company in selected_range.iterrows():
            
            attribute_list = []
            for attribute in company:
                print(attribute)
                if  (isinstance(attribute, int) or isinstance(attribute, str)):
                    print("not error")
                    attribute_list.append(attribute)
                else:
                    print("error")
                    attribute_list.append(0)
            
            company_temp = Company(attribute_list[0],
                            attribute_list[1],
                            self.all_employee[index],
                            self.all_quarter_performance[index],
                            attribute_list[2],
                            attribute_list[3],
                            attribute_list[4])
            self.all_companies.append(company_temp)

    def import_csv(self,  file):
        df_csv = pd.read_csv(file)
        client_List = []            # 1D array of Client objects
        for index, client in df_csv.iterrows():
            new_client = []
            for attribute in client:
                new_client.append(attribute)
            client_temp = Client(new_client[0],
                            new_client[1],
                            new_client[2],
                            new_client[3],
                            new_client[4],
                            new_client[5],
                            new_client[6])
            
            client_List.append(client_temp)
        self.all_companies[0].client_setter(client_List)

    def import_pdf(self, file):
        # The PDFPlumber library does not have a striaghtforward way to extract out specific text bt locale. rather it simply extracts all
        # text on a page as 1 string. So I had to slice the string to extarct out the page (table) header. Beyond that, the
        # table could be extracted simply
         
        with pdfplumber.open(file) as pdf:
            first_page = pdf.pages[0] 
            tables = first_page.extract_table()
            whole_page =  next(iter(first_page.extract_text_lines()))
            page_header = whole_page.get('text')            # Finds first occurance of 'text' key in the page, returns value
            df_pdf = pd.DataFrame(tables[1:], columns= tables [0])
        industry_performance_historical = []                # 2D array of quarters 
        
        for index, years in df_pdf.iterrows():
            quarter_Performance = []
            for attribute in years:
                try:
                    attribute = attribute.replace(",", "")
                except ValueError:
                    tmp = []
                except AttributeError:
                    tmp = []

                try: 
                    attribute = int(attribute)
                except ValueError:
                    tmp = []
                except AttributeError:
                    tmp = []

                quarter_Performance.append(attribute)

            industry_performance_historical.append(quarter_Performance)

        # Create Industry Object and insert all other sub-objects

        self.all_industry.append(Industry(page_header, industry_performance_historical, self.all_companies))


    def import_pptx(self, file):
        # Because the data may exist within different "shapes" such as text box, tables etc, they need to be parsed over and 
        # standardised into a common data structure: a dataframe. Due to the lack of a discernible pattern to the shape
        # order, I had to do some tailoring to the pptx format for each extracted shape
       
        prs = Presentation(file)
        extracted_df = []       # 1D array to store all shapes extracted from PDF

        # Check for tables or text frames
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    data = []

                    # Extract data rows
                    for row in table.rows:#[1:]
                        row_data = [cell.text for cell in row.cells]
                        data.append(row_data)

                    # Convert to DataFrame
                    if data:
                        df = pd.DataFrame(data[1:], columns=data[0])
                        extracted_df.append(df)

                # Check for text boxes
                elif shape.has_text_frame:
                    text_frame = shape.text_frame
                    extracted_text = []
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            extracted_text.append(run.text)
                    extracted_df.append(extracted_text)

            
        first_Shape = extracted_df[0]

        # Page 1 of PPTX
        report_year = int(extracted_df[0][0][-4:]) 
        total_revenue = int(extracted_df[1][1][16:].replace(",", ""))
        total_membership_sold = int(extracted_df[1][2][24:].replace(",", ""))
        top_location = extracted_df[1][3][14:]

        # Page 2 of PPTX
        quarter_performance_list = []
        
        for index, quarters in extracted_df[3].iterrows():
            quarter_performance_tmp_obj = QuarterPerformance( report_year,
                                                   quarters.iloc[0],
                                                   int(quarters.iloc[2]),
                                                   int(quarters.iloc[3]),
                                                   int(quarters.iloc[1].replace(",", "")))
            
            quarter_performance_list.append(quarter_performance_tmp_obj)         

        # Page 3 of PPTX
        distribution_list = []
        distribution_list.append(int(extracted_df[5][1][5:7]))
        distribution_list.append(int(extracted_df[5][2][6:8]))
        distribution_list.append(int(extracted_df[5][3][14:16]))
        distribution_list.append(int(extracted_df[5][4][19:21]))
        
        annual_temp = AnnualPerformance(report_year, 
                                    total_revenue,
                                    total_membership_sold,
                                    top_location,
                                    quarter_performance_list,
                                    distribution_list )
        self.all_industry[0].company_List[0].annual_performance_add(annual_temp)
    
    def Export (self):
        file = self.all_industry[0].get_dict()
        with open('unified_Data.json', 'w') as json_file:
            json.dump(file, json_file, indent=4)
        

#===========================================================================================================================

import_Temp = Importer("dataset1.json", "dataset2.csv", "dataset3.pdf", "dataset4.pptx")
import_Temp.Export()

